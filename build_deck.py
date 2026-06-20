# -*- coding: utf-8 -*-
"""Build the master's-thesis progress presentation (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PIL import Image

BASE = r"C:\Users\herma\OneDrive\Skrivebord\Claude work"

# ----- palette -----
INK     = RGBColor(0x12, 0x30, 0x28)   # deep spruce (dark slides)
INK2    = RGBColor(0x1B, 0x45, 0x39)   # lighter green panel
FOREST  = RGBColor(0x2C, 0x5F, 0x2D)
TEAL    = RGBColor(0x1C, 0x72, 0x93)   # accent (viridis teal)
AMBER   = RGBColor(0xC4, 0x86, 0x1A)   # accent 2 (low-sampling yellow)
MOSS    = RGBColor(0x6E, 0x9B, 0x5E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xEE, 0xF3, 0xEC)   # light green-gray card
PANEL2  = RGBColor(0xE7, 0xEF, 0xF1)   # light teal-gray card
TEXT    = RGBColor(0x1E, 0x2B, 0x25)
MUTED   = RGBColor(0x5E, 0x6E, 0x63)
LINE    = RGBColor(0xD3, 0xDD, 0xD0)
CREAMTX = RGBColor(0xD7, 0xE4, 0xDB)   # light text on dark

HEAD = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False  # inserts a single empty <a:effectLst/>
    if shadow:
        spPr = sp._element.spPr
        ef = spPr.find(qn('a:effectLst'))
        if ef is None:
            ef = spPr.makeelement(qn('a:effectLst'), {})
            spPr.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad':'60000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'1E2B25'})
        alpha = clr.makeelement(qn('a:alpha'), {'val':'22000'})
        clr.append(alpha); sh.append(clr); ef.append(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def _set_run(r, text, size, color, bold=False, italic=False, font=BODY, spacing=None):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        rpr = r._r.get_or_add_rPr()
        rpr.set('spc', str(int(spacing*100)))

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph = list of run dicts."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for rd in para:
            r = p.add_run()
            _set_run(r, rd.get("t",""), rd.get("sz",16), rd.get("c",TEXT),
                     rd.get("b",False), rd.get("i",False), rd.get("f",BODY), rd.get("sp"))
    return tb

def bullets(s, x, y, w, h, items, size=16, gap=9, bullet_color=TEAL, text_color=TEXT,
            line_spacing=1.04):
    """items: list of (text, level) or (text, level, color)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first = True
    for it in items:
        txt = it[0]; lvl = it[1] if len(it) > 1 else 0
        tcol = it[2] if len(it) > 2 else text_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else max(3, gap-4))
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        p.level = lvl
        glyph = "  " if lvl > 0 else ""
        indent = "      " * lvl
        rb = p.add_run()
        mark = "–  " if lvl > 0 else "▪  "
        _set_run(rb, indent+mark, size if lvl==0 else size-1,
                 AMBER if lvl>0 else bullet_color, bold=True, font=BODY)
        rt = p.add_run()
        _set_run(rt, txt, size if lvl==0 else size-1, tcol, font=BODY)
    return tb

def picture_fit(s, path, x, y, max_w, max_h, align="center", valign="top"):
    iw, ih = Image.open(path).size
    ratio = iw/ih
    w = max_w; h = w/ratio
    if h > max_h:
        h = max_h; w = h*ratio
    px = x + (max_w - w)/2 if align=="center" else (x + (max_w-w) if align=="right" else x)
    py = y + (max_h - h)/2 if valign=="middle" else (y + (max_h-h) if valign=="bottom" else y)
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
    return px, py, w, h

def eyebrow_title(s, eyebrow, title, num, dark=False):
    tc = WHITE if dark else INK
    text(s, 0.7, 0.40, 11.9, 0.3,
         [[{"t":eyebrow,"sz":12,"c":TEAL if not dark else MOSS,"b":True,"f":BODY,"sp":2.2}]])
    text(s, 0.7, 0.66, 11.0, 0.85,
         [[{"t":title,"sz":29,"c":tc,"b":True,"f":HEAD}]])
    # slide-number chip
    rect(s, SW-1.15, 0.46, 0.5, 0.5, fill=(TEAL if not dark else MOSS),
         shape=MSO_SHAPE.OVAL)
    text(s, SW-1.15, 0.46, 0.5, 0.5, [[{"t":str(num),"sz":15,"c":WHITE,"b":True,"f":HEAD}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def footer(s, label, dark=False):
    c = CREAMTX if dark else MUTED
    text(s, 0.7, 7.06, 8.0, 0.3,
         [[{"t":"Master's thesis  ·  Progress review", "sz":9, "c":c, "f":BODY}]])
    text(s, SW-4.7, 7.06, 4.0, 0.3,
         [[{"t":label, "sz":9, "c":c, "f":BODY}]], align=PP_ALIGN.RIGHT)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

P = lambda f: os.path.join(BASE, f)

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, 0.16, fill=AMBER)               # thin top accent
# left text block
text(s, 0.85, 1.35, 7.2, 0.3,
     [[{"t":"MASTER'S THESIS  ·  PROGRESS PRESENTATION","sz":13,"c":MOSS,"b":True,"f":BODY,"sp":2.5}]])
text(s, 0.85, 1.95, 7.4, 2.2,
     [[{"t":"Regional Structure of Sampling Bias in Species-Occurrence Data","sz":33,"c":WHITE,"b":True,"f":HEAD}]],
     line_spacing=1.04)
text(s, 0.85, 3.95, 7.2, 0.9,
     [[{"t":"Comparing modelled sampling-intensity surfaces across taxonomic groups in mainland Norway","sz":16,"c":CREAMTX,"i":True,"f":BODY}]],
     line_spacing=1.1)
# author / supervisor / date
text(s, 0.88, 5.15, 7.2, 1.0,
     [[{"t":"[ Author name ]", "sz":17, "c":WHITE, "b":True, "f":BODY}],
      [{"t":"Supervisor: [ to confirm ]", "sz":14, "c":CREAMTX, "f":BODY}],
      [{"t":"June 2026", "sz":14, "c":CREAMTX, "f":BODY}]],
     space_after=4)
# flag note
rect(s, 0.85, 6.55, 6.7, 0.55, fill=INK2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, 1.05, 6.55, 6.4, 0.55,
     [[{"t":"⚠  Fields in [brackets] need your input — see speaker notes.","sz":11.5,"c":AMBER,"b":True,"f":BODY}]],
     anchor=MSO_ANCHOR.MIDDLE)
# right framed map card
cx, cy, cw, ch = 8.55, 1.55, 4.15, 4.55
rect(s, cx, cy, cw, ch, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
picture_fit(s, P(r"birds\full_maps_png\birds1.png"), cx+0.2, cy+0.2, cw-0.4, ch-0.75, valign="top")
text(s, cx+0.2, cy+ch-0.5, cw-0.4, 0.4,
     [[{"t":"Modelled sampling intensity over Norway (one bird group)","sz":9.5,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "TITLE SLIDE — please confirm/replace before presenting:\n"
        "• [Author name] — your full name.\n"
        "• Supervisor — I could not find a supervisor named anywhere in the repo. Add their name (and co-supervisor if any).\n"
        "• Thesis title — the title shown is a WORKING/DESCRIPTIVE title I wrote from the code and data; replace it with your official registered thesis title if it differs.\n"
        "• Date — set to June 2026 (today's date context). Adjust to your actual presentation date.\n\n"
        "Talking point: 'This is a progress update on my thesis, which looks at how sampling effort — where people actually record species — is spatially biased across Norway, and whether that bias pattern is shared across different taxonomic groups. The map on the right is the central object I work with: a modelled sampling-intensity surface.'")

# =====================================================================
# SLIDE 2 — WHY IT MATTERS
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "BACKGROUND  ·  MOTIVATION", "Why sampling bias matters", 2)
bullets(s, 0.75, 1.75, 6.7, 4.6, [
    ("Biodiversity records (GBIF, citizen science, museums) are gathered opportunistically — not by survey design.", 0),
    ("Where people look ≠ where species live. Roads, towns and popular areas are over-recorded; remote areas under-recorded.", 0),
    ("This spatial sampling bias leaks into everything built on the data:", 0),
    ("species distribution models", 1),
    ("richness and rarity maps", 1),
    ("conservation and monitoring decisions", 1),
    ("Modern workflows estimate a per-group sampling-intensity (“bias”) surface so the effort can be modelled and corrected.", 0),
], size=16, gap=10)
# right illustrative map in a card
cx, cy, cw, ch = 8.0, 1.7, 4.6, 4.75
rect(s, cx, cy, cw, ch, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
px,py,w,h = picture_fit(s, P(r"maps\turdus_pilaris_norway.png"), cx+0.25, cy+0.25, cw-0.5, ch-0.95, valign="top")
text(s, cx+0.25, cy+ch-0.62, cw-0.5, 0.5,
     [[{"t":"Fieldfare (Turdus pilaris): 484,709 records — dense in the south & along the coast, sparse in the interior north.",
        "sz":10,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Why it matters")
notes(s, "Set up the problem for a general audience.\n\n"
        "'Almost all large biodiversity datasets are opportunistic — they record where observers happened to be, not a designed sample of the country. So the data tells you as much about human behaviour as about the species.'\n\n"
        "'If you don't account for that, a distribution model can mistake well-surveyed for suitable habitat. The standard fix is to estimate a sampling-intensity or bias surface for a group of similar species and use it to correct the effort.'\n\n"
        "Point at the fieldfare map: even for an abundant, easy-to-see bird with nearly half a million records, you can see the classic bias — dense south and coast, thin in the northern interior. My thesis asks how consistent that bias pattern is across very different groups.")

# =====================================================================
# SLIDE 3 — RESEARCH QUESTION
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "THE GAP  ·  AIMS", "Research question", 3)
text(s, 0.75, 1.6, 11.8, 0.95,
     [[{"t":"Is the spatial structure of sampling bias ","sz":21,"c":TEXT,"f":HEAD},
       {"t":"shared across taxonomic groups","sz":21,"c":TEAL,"b":True,"f":HEAD},
       {"t":", or is it group- and region-specific?","sz":21,"c":TEXT,"f":HEAD}]],
     line_spacing=1.1)
# three question cards
cards = [
    ("01", "Aggregate", "Across many groups, is the overall spread of sampling intensity a stable property, or do particular groups dominate it?"),
    ("02", "Spatial", "Which regions drive differences between groups? Are some regions interchangeable and others decisive?"),
    ("03", "Why it matters", "If bias is shared, one correction can serve many groups. If it is group-specific, corrections must be tailored."),
]
cw = 3.86; gap = 0.27; x0 = 0.75; y0 = 2.95; ch = 3.05
for i,(num,hd,body) in enumerate(cards):
    x = x0 + i*(cw+gap)
    rect(s, x, y0, cw, ch, fill=PANEL if i!=2 else PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    rect(s, x, y0, 0.12, ch, fill=[TEAL,FOREST,AMBER][i])
    text(s, x+0.35, y0+0.28, cw-0.6, 0.6, [[{"t":num,"sz":26,"c":[TEAL,FOREST,AMBER][i],"b":True,"f":HEAD}]])
    text(s, x+0.35, y0+0.95, cw-0.6, 0.5, [[{"t":hd,"sz":17,"c":INK,"b":True,"f":HEAD}]])
    text(s, x+0.35, y0+1.5, cw-0.62, 1.4, [[{"t":body,"sz":13.5,"c":TEXT,"f":BODY}]], line_spacing=1.08)
footer(s, "Research question")
notes(s, "State the question plainly, then unpack the three angles on the cards.\n\n"
        "Core question: is sampling bias a shared backdrop you can correct once, or does every group and region behave differently?\n\n"
        "Card 1 (Aggregate): pool all the cell values across groups and ask whether the total variance is something stable, or whether swapping which groups you include moves it a lot. This is the permutation test later.\n\n"
        "Card 2 (Spatial): break it down by region — which regions, when you change the group they come from, actually move the result? This is the block-randomization / region-sensitivity test.\n\n"
        "Card 3 (Stakes): the practical payoff — a shared bias structure means a generic effort correction is defensible; a group-specific one means you must model effort per group.")

# =====================================================================
# SLIDE 4 — DATA & STUDY SYSTEM
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "DATA  ·  STUDY SYSTEM", "Data & study system", 4)
# two contrasting maps
mw = 3.05
picture_fit(s, P(r"maps\turdus_pilaris_norway.png"), 0.7, 1.7, mw, 3.5, valign="top")
picture_fit(s, P(r"maps\fomitopsis_pinicola_norway.png"), 0.7+mw+0.15, 1.7, mw, 3.5, valign="top")
text(s, 0.7, 5.25, mw, 0.4, [[{"t":"Fieldfare — 484,709","sz":11,"c":MUTED,"b":True,"f":BODY}]], align=PP_ALIGN.CENTER)
text(s, 0.7+mw+0.15, 5.25, mw, 0.4, [[{"t":"Red-belted conk — 9,585","sz":11,"c":MUTED,"b":True,"f":BODY}]], align=PP_ALIGN.CENTER)
text(s, 0.7, 5.6, mw*2+0.15, 0.4,
     [[{"t":"Same country, very different coverage — bias is group-dependent.","sz":11,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER)
# right column text
rx = 7.35
bullets(s, rx, 1.7, 5.3, 2.6, [
    ("Three taxon datasets of modelled groups: birds (21), fungi (20), vascular plants (19), plus a bird re-run “newbirds” (20).", 0),
    ("Each group has a posterior sampling-intensity (bias-mean) raster over mainland Norway, UTM33N, 10 km grid.", 0),
    ("Four exemplar GBIF species illustrate the raw effort behind the surfaces.", 0),
], size=14.5, gap=9)
# small species count table
tbl = [
    ("Turdus pilaris  (bird)", "484,709"),
    ("Falco peregrinus  (bird)", "79,445"),
    ("Lysimachia europaea  (plant)", "29,605"),
    ("Fomitopsis pinicola  (fungus)", "9,585"),
]
ty = 4.45
rect(s, rx, ty, 5.3, 1.55, fill=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, rx+0.25, ty+0.12, 5.0, 0.3, [[{"t":"GBIF occurrences, mainland Norway  ·  DOI 10.15468/dl.3ru6rq","sz":10.5,"c":TEAL,"b":True,"f":BODY}]])
for i,(nm,ct) in enumerate(tbl):
    yy = ty+0.45+i*0.27
    text(s, rx+0.25, yy, 3.6, 0.27, [[{"t":nm,"sz":12,"c":TEXT,"i":True,"f":BODY}]])
    text(s, rx+0.25, yy, 4.8, 0.27, [[{"t":ct,"sz":12,"c":INK,"b":True,"f":BODY}]], align=PP_ALIGN.RIGHT)
footer(s, "Data & study system")
notes(s, "Describe the data without drowning the audience.\n\n"
        "'I work with modelled sampling-intensity surfaces for three taxonomic datasets — birds, fungi and vascular plants — each split into ~20 groups, plus a re-run of the birds I call newbirds. Each group gives one raster: the posterior mean of estimated sampling effort over mainland Norway on a 10 km grid.'\n\n"
        "'To ground that in something concrete, I pulled four real species straight from GBIF as a formal, citable download (one DOI). The counts span two orders of magnitude — half a million fieldfare records versus under ten thousand for the red-belted conk — which already shows the bias is strongly group-dependent.'\n\n"
        "CAVEAT to mention if asked: the exact definition of a 'group' (functional vs taxonomic split, or model replicate) comes from the upstream modelling step — confirm the precise wording before the committee asks.")

# =====================================================================
# SLIDE 5 — THE CORE OBJECT: BIAS SURFACES
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "WHAT I ANALYSE", "From records to a sampling-intensity surface", 5)
picture_fit(s, P(r"birds\full_maps_png\birds1.png"), 0.7, 1.65, 4.35, 4.95, valign="top")
rx = 5.5
bullets(s, rx, 1.75, 7.1, 4.4, [
    ("Each group's raw records are turned into a smooth sampling-intensity surface (log scale) by the upstream model.", 0),
    ("High values = intensively recorded; low values = under-sampled. The north-interior is consistently the thinnest.", 0),
    ("These surfaces — not the raw points — are the unit of analysis: they are directly comparable across groups on one grid.", 0),
    ("My contribution starts here: quantifying how the surfaces vary, and where, across all groups within a dataset.", 0),
], size=16, gap=11)
rect(s, rx, 5.55, 7.1, 0.95, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, rx+0.3, 5.55, 6.6, 0.95,
     [[{"t":"Key idea:  ","sz":14,"c":TEAL,"b":True,"f":BODY},
       {"t":"one comparable surface per group lets me test whether the bias pattern is a shared property or a group-by-region one.","sz":14,"c":TEXT,"f":BODY}]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, "Sampling-intensity surfaces")
notes(s, "This slide is the hinge between background and my own work.\n\n"
        "'The modelling step converts each group's scattered records into a continuous surface of estimated sampling intensity on a common 10 km grid. Yellow-green is well recorded, blue-purple is poorly recorded. Notice the persistent low-intensity zone in the northern interior — that recurs across groups.'\n\n"
        "'Crucially, I analyse these surfaces rather than the raw points, because surfaces for different groups live on the same grid and are directly comparable. Everything from here on — the regional sampling, the variance tests — is my own downstream analysis of these surfaces.'")

# =====================================================================
# SLIDE 6 — REGIONAL SAMPLING DESIGN
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "METHOD  ·  REGIONAL DESIGN", "Sampling 12 regions across Norway", 6)
picture_fit(s, P(r"birds\region_maps_png\birds1.png"), 0.7, 1.7, 7.35, 4.55, valign="top")
text(s, 0.7, 6.25, 7.35, 0.4,
     [[{"t":"Full surface with 12 region windows, plus the zoomed 10×10 km blocks.","sz":10,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER)
rx = 8.25
bullets(s, rx, 1.72, 4.4, 4.5, [
    ("12 fixed regions spanning south-to-north and coast-to-interior (e.g. Oslo, Bergen, Trondheim, Tromsø, Kirkenes).", 0),
    ("Each region = a 10×10 km window, giving 121 grid cells of bias values per group.", 0),
    ("Started with 6 regions, extended to 12 for fuller north coverage.", 0),
    ("Refinement: Svolvær & Kirkenes centres nudged a few km so every block holds a full 121 land cells (no ocean/NA) — keeping regions strictly comparable.", 0),
], size=14.5, gap=10)
footer(s, "Regional design")
notes(s, "Explain the sampling scheme and a bit of the care taken.\n\n"
        "'To compare groups fairly I don't use the whole country at once — I sample 12 fixed regions chosen to span the gradients that matter: south to north, coast to interior. Each region is a 10-by-10 km window, which is exactly 121 cells of the bias surface for every group.'\n\n"
        "'I started with six regions and doubled it to twelve to get better northern coverage. One practical detail: two northern windows — Svolvær and Kirkenes — originally clipped the coastline and lost cells to ocean/NA. I shifted their centres by a few kilometres so all twelve regions have a complete 121 land cells across every group. That's what makes the variance comparisons clean.'")

# =====================================================================
# SLIDE 7 — RESULT 1: PERMUTATION TEST
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "RESULTS  ·  1 of 2", "Total spread of bias is a stable property", 7)
picture_fit(s, P("birds_permutation_histogram.png"), 0.7, 1.7, 6.55, 4.2, valign="top")
text(s, 0.7, 5.85, 6.55, 0.4,
     [[{"t":"Birds: observed total variance (red) sits in the centre of the null.","sz":10,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER)
rx = 7.55
# table
rows = [
    ["Dataset","Grp","Observed","Null mean","P(≥obs)"],
    ["Birds","21","0.2405","0.2402","0.477"],
    ["Fungi","20","0.4653","0.4652","0.478"],
    ["Vasc. plants","19","0.4807","0.4796","0.479"],
]
tx, ty, tw = rx, 1.75, 5.1
colw = [1.55, 0.62, 1.05, 1.05, 0.83]
rh = 0.46
for ri,row in enumerate(rows):
    cx = tx
    for ci,cell in enumerate(row):
        is_head = ri==0
        if is_head:
            rect(s, cx, ty+ri*rh, colw[ci], rh, fill=INK)
        elif ri%2==0:
            rect(s, cx, ty+ri*rh, colw[ci], rh, fill=PANEL)
        col = WHITE if is_head else TEXT
        b = is_head or ci==0
        al = PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        text(s, cx+(0.12 if ci==0 else 0), ty+ri*rh, colw[ci]-(0.12 if ci==0 else 0), rh,
             [[{"t":cell,"sz":12.5,"c":col,"b":b,"f":BODY}]], align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
# callout
co_y = ty + len(rows)*rh + 0.35
rect(s, rx, co_y, tw, 1.95, fill=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
rect(s, rx, co_y, 0.12, 1.95, fill=TEAL)
text(s, rx+0.32, co_y+0.18, tw-0.55, 0.4, [[{"t":"What the test does","sz":13,"c":TEAL,"b":True,"f":BODY}]])
text(s, rx+0.32, co_y+0.55, tw-0.55, 1.35,
     [[{"t":"1,000 permutations re-draw groups (with replacement) per region and recompute the pooled variance. Across all three datasets the observed value lands mid-distribution (≈ 48% of draws exceed it).",
        "sz":12.5,"c":TEXT,"f":BODY}],
      [{"t":"→ No single group dominates; the total spread of sampling intensity is exchangeable across groups.",
        "sz":12.5,"c":INK,"b":True,"f":BODY}]],
     line_spacing=1.06, space_after=6)
footer(s, "Result 1 — permutation test")
notes(s, "First main result. Keep the takeaway crisp.\n\n"
        "'I pool all the region-by-group cell values and compute their total variance, with each group used once — that's the observed baseline. Then I run a permutation: for each region I redraw groups with replacement and recompute the variance, a thousand times, to get a null distribution.'\n\n"
        "'The result is the same for all three datasets: the observed value sits right in the middle of the null — about 48% of permutations exceed it. In other words, the total spread of sampling intensity doesn't depend on which particular groups you include. No group is an outlier driving the variance.'\n\n"
        "Interpretation to offer: this is a meaningful negative/stability result — at the aggregate level, bias intensity is a shared, exchangeable property. The interesting structure shows up only when you break it down spatially (next slide).\n\n"
        "Numbers are on a log scale; seed fixed at 1234 for reproducibility.")

# =====================================================================
# SLIDE 8 — RESULT 2: REGION SENSITIVITY
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "RESULTS  ·  2 of 2", "Northern regions drive group differences", 8)
# bar chart (horizontal). Categories bottom-to-top -> list ascending so largest on top.
regions = ["Oslo","Valdres","Svolvær","Trondheim","Bergen","Kristiansand","Bodø","Setesdal","Tromsø","Lakselv","Kirkenes","Skorovatn"]
vals    = [0.14,0.53,0.58,0.68,0.70,0.73,0.75,0.77,1.98,2.97,18.12,22.12]
cd = CategoryChartData()
cd.categories = regions
cd.add_series("Mean abs. change", vals)
gx, gy, gw, gh = 0.7, 1.7, 6.7, 4.55
chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(gx), Inches(gy), Inches(gw), Inches(gh), cd).chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0'; plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(9); plot.data_labels.font.color.rgb = MUTED
series = plot.series[0]
series.format.fill.solid(); series.format.fill.fore_color.rgb = TEAL
# color the two dominant bars amber via point-level fill
from pptx.oxml.ns import qn as _qn
for idx in (10,11):
    pt = series.points[idx]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = AMBER
cat_ax = chart.category_axis; val_ax = chart.value_axis
cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.name = BODY
val_ax.tick_labels.font.size = Pt(9)
val_ax.has_major_gridlines = True
cat_ax.has_major_gridlines = False
try:
    val_ax.axis_title.text_frame.text = ""
except Exception:
    pass
text(s, gx, gy+gh+0.02, gw, 0.35,
     [[{"t":"Mean absolute change in variance-among-regions when each region is swapped  (×10⁻⁶).","sz":10,"c":MUTED,"i":True,"f":BODY}]],
     align=PP_ALIGN.CENTER)
rx = 7.7
bullets(s, rx, 1.72, 4.95, 3.0, [
    ("A second test swaps one region's data from a donor group and measures how much the between-region variance moves.", 0),
    ("Two northern, data-sparse regions — Skorovatn and Kirkenes — dominate; swapping them moves the result 10×+ more than any southern region.", 0),
    ("Well-sampled southern regions (Oslo, Bergen, Valdres) are effectively interchangeable across groups.", 0),
], size=14, gap=9)
rect(s, rx, 4.95, 4.95, 1.5, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
rect(s, rx, 4.95, 0.12, 1.5, fill=AMBER)
text(s, rx+0.3, 5.08, 4.5, 1.3,
     [[{"t":"Where bias differs between groups is spatially concentrated","sz":13,"c":AMBER,"b":True,"f":BODY}],
      [{"t":"— in the under-recorded north, not the well-surveyed south.","sz":13,"c":TEXT,"f":BODY}]],
     line_spacing=1.05, space_after=3)
footer(s, "Result 2 — region sensitivity")
notes(s, "Second main result — the spatial counterpart to slide 7.\n\n"
        "'The first test said the aggregate is stable. This one asks WHERE the group-to-group differences live. The statistic is the variance among the 12 regional variances. I swap one region at a time, taking that region's cells from a different (donor) group, and see how much the statistic moves. Averaging over all donors gives a sensitivity score per region.'\n\n"
        "'The answer is sharply spatial: Skorovatn and Kirkenes — both northern and data-poor — dominate. Swapping them changes the result more than ten times as much as any southern region. Oslo and the other well-sampled southern regions barely move it; they're essentially interchangeable across groups.'\n\n"
        "IMPORTANT honesty note: these sensitivity numbers are currently from the 'newbirds' (bird re-run) dataset, the 12-region pipeline output. Running the same region-sensitivity breakdown for fungi and vascular plants is part of the remaining work (see future-work slide).\n\n"
        "Ecological reading to float: the north is where effort is thinnest and most variable between groups, so it carries the most leverage — which is exactly where a generic bias correction is riskiest.")

# =====================================================================
# SLIDE 9 — SYNTHESIS
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "INTERPRETATION", "What the results say so far", 9)
items = [
    ("At the aggregate level, bias is shared", "The total spread of sampling intensity is exchangeable across groups — no group is special. A common baseline is defensible.", TEAL),
    ("Spatially, the differences are concentrated", "Group-to-group variation lives in a few northern, under-sampled regions; the well-surveyed south is stable.", AMBER),
    ("Implication for correction", "A one-size correction is reasonable in the south but risky in the data-sparse north, where group-specific effort matters most.", FOREST),
]
y0 = 1.85; ch = 1.42; gap = 0.22
for i,(hd,body,col) in enumerate(items):
    y = y0 + i*(ch+gap)
    rect(s, 0.75, y, 11.85, ch, fill=PANEL if i%2==0 else PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, 0.75, y, 0.14, ch, fill=col)
    rect(s, 1.15, y+0.34, 0.74, 0.74, fill=col, shape=MSO_SHAPE.OVAL)
    text(s, 1.15, y+0.34, 0.74, 0.74, [[{"t":str(i+1),"sz":24,"c":WHITE,"b":True,"f":HEAD}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.15, y+0.22, 10.2, 0.5, [[{"t":hd,"sz":17,"c":INK,"b":True,"f":HEAD}]])
    text(s, 2.15, y+0.72, 10.2, 0.6, [[{"t":body,"sz":14,"c":TEXT,"f":BODY}]], line_spacing=1.04)
footer(s, "Synthesis")
notes(s, "Pull the two results into one story.\n\n"
        "'Putting the two tests together: zoom out and bias looks like a shared backdrop — the overall variance doesn't care which groups you include. Zoom in and the story changes — the differences between groups are concentrated in a handful of northern, poorly-sampled regions.'\n\n"
        "'The practical message for anyone correcting sampling bias: in the well-surveyed south you can largely get away with a shared correction; in the data-sparse north, group-specific effort genuinely matters and a generic correction is most likely to mislead.'\n\n"
        "Flag that these are progress-stage conclusions: directionally consistent, but significance and cross-taxon replication still to be finalised.")

# =====================================================================
# SLIDE 10 — IMPLEMENTATION
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "UNDER THE HOOD", "How it was built", 10)
# pipeline flow
steps = ["Config","Setup","Full maps","Region maps","Extract","Randomize","Variance test"]
n = len(steps); pw = 1.52; ph = 0.78; gap = 0.18; x0 = 0.75; y0 = 1.7
for i,st in enumerate(steps):
    x = x0 + i*(pw+gap)
    col = INK if i in (0,) else (AMBER if i==n-1 else TEAL)
    rect(s, x, y0, pw, ph, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    text(s, x, y0, pw, ph, [[{"t":f"{i:02d}","sz":10,"c":CREAMTX,"b":True,"f":BODY}],
                            [{"t":st,"sz":11.5,"c":WHITE,"b":True,"f":BODY}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    if i < n-1:
        text(s, x+pw-0.04, y0, gap+0.08, ph, [[{"t":"›","sz":18,"c":MUTED,"b":True,"f":BODY}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.75, 2.95, 6.1, 3.4, [
    ("R + terra / sf for raster & spatial work; rnaturalearth for the Norway outline.", 0),
    ("Modular pipeline driven by one config file — switching datasets edits a single script.", 0),
    ("Bias surfaces read as posterior-mean layers, masked to Norway on a common grid.", 0),
    ("Permutation & variance tests written separately and run read-only over the data.", 0),
], size=14.5, gap=10)
# reproducibility card
rx = 7.15
rect(s, rx, 2.95, 5.45, 3.25, fill=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
text(s, rx+0.32, 3.12, 4.9, 0.4, [[{"t":"Reproducibility by design","sz":14,"c":TEAL,"b":True,"f":BODY}]])
bullets(s, rx+0.32, 3.6, 4.85, 2.5, [
    ("Formal GBIF download via occ_download() → one citable DOI (10.15468/dl.3ru6rq).", 0),
    ("Credentials read from .Renviron — no secrets in code.", 0),
    ("Fixed random seed (1234) for every permutation run.", 0),
    ("Region centres validated to a full 121 land cells across all groups.", 0),
], size=13, gap=8, bullet_color=TEAL)
footer(s, "Implementation")
notes(s, "Keep this understandable, not a code walkthrough.\n\n"
        "'Everything is in R, built around the terra and sf spatial stack. The plotting and extraction live in a modular pipeline — config, setup, full maps, region maps, value extraction, randomization — so moving from birds to fungi to plants is just editing one config file.'\n\n"
        "'I cared about reproducibility: the example occurrences come from a formal GBIF download that produces a single citable DOI, credentials are kept out of the code in an environment file, every permutation uses a fixed seed, and the region windows are validated to a full 121 land cells so the comparisons are exact.'\n\n"
        "'The statistical tests — the permutation and the randomization — are written as separate, strictly read-only scripts over the data, so analysis never mutates the modelling outputs.'\n\n"
        "If asked about the upstream model that produces the bias surfaces: that's an existing Bayesian sampling-effort model; my work is the downstream comparative analysis of its outputs.")

# =====================================================================
# SLIDE 11 — FUTURE WORK & TIMELINE
# =====================================================================
s = slide(); bg(s, WHITE)
eyebrow_title(s, "WHAT'S LEFT", "Remaining work & timeline", 11)
bullets(s, 0.75, 1.7, 6.2, 4.5, [
    ("Run the 12-region permutation AND region-sensitivity tests consistently across all taxa (finalise newbirds; add fungi & plants to Result 2).", 0),
    ("Unify the two variance statistics — pooled total variance and variance-among-regions — into one framework.", 0),
    ("Add proper significance & sensitivity to the choice and number of regions.", 0),
    ("Interpret the ecological drivers of the northern-region leverage (accessibility, effort, habitat).", 0),
    ("Write up: methods, results and discussion chapters.", 0),
], size=14.5, gap=11)
# timeline card (flagged as proposed)
rx = 7.3
rect(s, rx, 1.7, 5.33, 4.5, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
text(s, rx+0.32, 1.88, 4.7, 0.35, [[{"t":"PROPOSED TIMELINE","sz":12,"c":AMBER,"b":True,"f":BODY,"sp":1.5}]])
text(s, rx+0.32, 2.2, 4.7, 0.3, [[{"t":"⚠  No timeline found in the repo — confirm/replace.","sz":10.5,"c":MUTED,"i":True,"f":BODY}]])
tl = [
    ("Now → summer", "Finish cross-taxon analysis; lock methods"),
    ("Late summer", "Significance + region-robustness checks"),
    ("Autumn", "Draft results & discussion"),
    ("Before submission", "Full write-up, revision, defence prep"),
]
yy = 2.7
for i,(when,what) in enumerate(tl):
    rect(s, rx+0.4, yy+0.05, 0.18, 0.18, fill=TEAL, shape=MSO_SHAPE.OVAL)
    if i < len(tl)-1:
        rect(s, rx+0.475, yy+0.23, 0.03, 0.62, fill=LINE)
    text(s, rx+0.75, yy-0.04, 4.2, 0.3, [[{"t":when,"sz":13,"c":INK,"b":True,"f":BODY}]])
    text(s, rx+0.75, yy+0.27, 4.2, 0.4, [[{"t":what,"sz":12,"c":TEXT,"f":BODY}]], line_spacing=1.0)
    yy += 0.85
footer(s, "Future work & timeline")
notes(s, "Set out what remains and be explicit that the timeline needs your input.\n\n"
        "'The main remaining task is consistency: the aggregate permutation test is done for all three taxa, but the region-sensitivity breakdown is currently only on the bird re-run. I need to run that for fungi and plants too, then reconcile the two variance statistics into a single framework and add proper significance testing and robustness to how the regions are chosen.'\n\n"
        "'After that it's interpretation — why the north carries the leverage — and writing the methods, results and discussion chapters.'\n\n"
        "TIMELINE — IMPORTANT: I could not find any dated plan or milestone notes in the repo, so the timeline on the right is a PLAUSIBLE SKELETON I drafted, not your real schedule. Replace it with your actual milestones and submission date before presenting, and check it against what your supervisor expects.")

# =====================================================================
# SLIDE 12 — CLOSING
# =====================================================================
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, 0.16, fill=AMBER)
text(s, 0.85, 1.05, 11.6, 0.35,
     [[{"t":"SUMMARY  ·  PROGRESS TO DATE","sz":13,"c":MOSS,"b":True,"f":BODY,"sp":2.5}]])
text(s, 0.85, 1.5, 11.6, 1.0,
     [[{"t":"Bias is a shared backdrop — except where the data runs thin","sz":29,"c":WHITE,"b":True,"f":HEAD}]],
     line_spacing=1.05)
pts = [
    ("Built", "A reproducible, multi-taxon pipeline turning GBIF records into comparable sampling-intensity surfaces over Norway."),
    ("Found", "Aggregate bias variance is exchangeable across groups; group differences concentrate in a few northern regions."),
    ("Next", "Extend Result 2 across all taxa, unify the variance framework, add significance, and write up."),
]
y0 = 3.05; cw = 3.86; gap = 0.27
for i,(hd,body) in enumerate(pts):
    x = 0.85 + i*(cw+gap)
    rect(s, x, y0, cw, 2.55, fill=INK2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    rect(s, x, y0, cw, 0.12, fill=[MOSS,AMBER,TEAL][i])
    hd_col = [MOSS, AMBER, RGBColor(0x7F, 0xC4, 0xDC)][i]
    text(s, x+0.32, y0+0.32, cw-0.6, 0.5, [[{"t":hd,"sz":18,"c":hd_col,"b":True,"f":HEAD}]])
    text(s, x+0.32, y0+0.92, cw-0.62, 1.5, [[{"t":body,"sz":13.5,"c":CREAMTX,"f":BODY}]], line_spacing=1.1)
text(s, 0.85, 6.05, 11.6, 0.5,
     [[{"t":"Thank you — questions welcome.","sz":15,"c":WHITE,"i":True,"f":HEAD}]])
notes(s, "Close confidently and invite discussion.\n\n"
        "'To sum up: I've built a reproducible pipeline that turns raw biodiversity records into comparable sampling-intensity surfaces across birds, fungi and plants. The headline finding is two-sided — at the whole-country scale, sampling bias behaves like a shared backdrop that doesn't depend on the group, but the group-to-group differences are real and they concentrate in the under-sampled north.'\n\n"
        "'Next I'll extend the spatial test across all taxa, tie the variance measures together, add significance, and move into writing. Happy to take questions.'\n\n"
        "Reminder: fill in your name, supervisor, official title (slide 1) and the real timeline (slide 11) before presenting.")

out = os.path.join(BASE, "Thesis_Progress_Presentation.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
